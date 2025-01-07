import streamlit as st
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
from io import BytesIO
import json
import time
from openai import AzureOpenAI
import openpyxl

# Set up your OpenAI API key
api_key = st.secrets["AZURE_OPENAI_API_KEY"]
azure_endpoint = st.secrets["AZURE_OPENAI_ENDPOINT"]

client = AzureOpenAI(
    api_key=api_key,  
    api_version="2024-02-01",
    azure_endpoint=azure_endpoint
)

def extract_all_text_from_ppt(file):
    ppt = Presentation(file)
    texts = {}
    full_text = []

    def extract_text_from_shapes(shapes, slide_index, path, texts, full_text):
        for index, shape in enumerate(shapes):
            current_path = f"{path},{index}"
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_text = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    shape_text.append(paragraph_text)
                    full_text.append(paragraph_text)
                if shape_text:
                    texts[current_path] = shape_text
            elif hasattr(shape, "shapes"):  # This is a group shape.
                extract_text_from_shapes(shape.shapes, slide_index, current_path, texts, full_text)

    for i, slide in enumerate(ppt.slides):
        extract_text_from_shapes(slide.shapes, i, f"{i}", texts, full_text)

    full_context = "\n\n".join(full_text)
    return texts, full_context

def extract_text_from_slide(slide):
    texts = {}
    full_text = []

    def extract_text_from_shapes(shapes, slide_index, path, texts, full_text):
        for index, shape in enumerate(shapes):
            current_path = f"{path},{index}"
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_text = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    shape_text.append(paragraph_text)
                    full_text.append(paragraph_text)
                if shape_text:
                    texts[current_path] = shape_text
            elif hasattr(shape, "shapes"):  # This is a group shape.
                extract_text_from_shapes(shape.shapes, slide_index, current_path, texts, full_text)

    extract_text_from_shapes(slide.shapes, slide.slide_id, f"{slide.slide_id}", texts, full_text)
    full_context = "\n\n".join(full_text)
    return texts, full_context

def apply_translated_text_to_slide(slide, translated_dict):
    def set_text_to_shape(shape, translated_paragraphs):
        if hasattr(shape, "text_frame"):
            for paragraph_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                if paragraph_idx < len(translated_paragraphs):
                    if paragraph.runs:
                        # Store the formatting of the first run
                        original_run = paragraph.runs[0]
                        font_size = original_run.font.size
                        font_color = original_run.font.color
                        bold = original_run.font.bold
                        italic = original_run.font.italic
                        underline = original_run.font.underline

                        # Clear existing text
                        paragraph.clear()

                        # Add new run with the translated text
                        new_run = paragraph.add_run()
                        new_run.text = translated_paragraphs[paragraph_idx]

                        # Apply the original formatting to the new run
                        if font_size:
                            new_run.font.size = font_size
                        if font_color and hasattr(font_color, 'rgb'):
                            new_run.font.color.rgb = font_color.rgb
                        new_run.font.bold = bold
                        new_run.font.italic = italic
                        new_run.font.underline = underline

    def apply_to_shapes(shapes, path):
        for index, shape in enumerate(shapes):
            current_path = f"{path},{index}"
            if current_path in translated_dict:  # Corrected from 'current_data' to 'current_path'
                set_text_to_shape(shape, translated_dict[current_path])
            elif hasattr(shape, "shapes"):
                apply_to_shapes(shape.shapes, current_path)  # Corrected recursive call

    apply_to_shapes(slide.shapes, f"{slide.slide_id}")

def translate_text(text_dict, target_language):
    max_retries = 2
    timeout_seconds = 120
    attempt = 0

    # system prompt:
    prompt = f"""
        You are a professional language translator.\n
        Return a json with format similar to user's provided dictionary\n
        Translate the text to {target_language}."""
    print("------------------- System Prompt:\n", prompt)

    # user prompt:
    converted_dict = json.dumps({str(k): v for k, v in text_dict.items()})
    print("------------------- Input text dictionary:\n", converted_dict)

    while attempt < max_retries:
        try:
            start_time = time.time()
            response = client.chat.completions.create(
                model="gpt-4o",
                response_format={"type": "json_object"},
                messages=[{"role": "system", "content": prompt},
                          {"role": "user", "content": converted_dict}],
                temperature=0.5,
                max_tokens=4000,
            )
            elapsed_time = time.time() - start_time

            if elapsed_time > timeout_seconds:
                raise TimeoutError("API call exceeded the time limit of 120 seconds.")

            print("------------------- Raw response from API:\n", response) 

            # Assuming the response format matches the input dictionary order
            translated_dict = json.loads(response.choices[0].message.content)
            return translated_dict

        except (json.JSONDecodeError, TimeoutError) as e:
            print("Attempt", attempt + 1, "failed:", str(e))
            attempt += 1
            if attempt == max_retries:
                st.error(f"Translation failed after {max_retries} attempts.")
                return text_dict
        except Exception as e:
            st.error(f"Error in translation: {str(e)}")
            return text_dict
        time.sleep(2)  # wait 2 seconds before retrying (if needed)

    return text_dict

def process_pptx(file, target_language, progress_bar, total_slides):
    ppt = Presentation(file)
    for i, slide in enumerate(ppt.slides):
        slide_text_dict, _ = extract_text_from_slide(slide)
        translated_slide_dict = translate_text(slide_text_dict, target_language)
        apply_translated_text_to_slide(slide, translated_slide_dict)
        # Update progress bar based on the slide index, value between 0.0 and 1.0
        current_progress = (i + 1) / total_slides
        progress_bar.progress(current_progress, text= f"Processing slide {i}/{total_slides}")
    return ppt

def save_pptx(ppt, original_file_name, language):
    output = BytesIO()
    ppt.save(output)
    output.seek(0)
    return output, f"{language}_translated_{original_file_name}"

def extract_text_from_docx(doc):
    texts = {}
    full_text = []
    # Extract text from paragraphs
    for para_index, paragraph in enumerate(doc.paragraphs):
        for run_index, run in enumerate(paragraph.runs):
            path = f"paragraph_{para_index},run_{run_index}"
            run_text = run.text
            texts[path] = [run_text]
            full_text.append(run_text)
    
    # Extract text from tables
    for table_index, table in enumerate(doc.tables):
        for row_index, row in enumerate(table.rows):
            for cell_index, cell in enumerate(row.cells):
                path = f"table_{table_index},row_{row_index},cell_{cell_index}"
                cell_text = cell.text
                if cell_text:
                    texts[path] = [cell_text]
                    full_text.append(cell_text)
    
    return texts, "\n\n".join(full_text)


def apply_translated_text_to_docx(doc, translated_dict):
    # Apply translated text to paragraphs
    for para_index, paragraph in enumerate(doc.paragraphs):
        for run_index, run in enumerate(paragraph.runs):
            path = f"paragraph_{para_index},run_{run_index}"
            if path in translated_dict:
                translated_text = translated_dict[path][0]
                run.text = translated_text

    # Apply translated text to tables
    for table_index, table in enumerate(doc.tables):
        for row_index, row in enumerate(table.rows):
            for cell_index, cell in enumerate(row.cells):
                path = f"table_{table_index},row_{row_index},cell_{cell_index}"
                if path in translated_dict:
                    cell.text = translated_dict[path][0]

def process_docx(file, target_language):
    doc = Document(file)
    text_dict, _ = extract_text_from_docx(doc)
    translated_text_dict = translate_text(text_dict, target_language)
    apply_translated_text_to_docx(doc, translated_text_dict)
    return doc

def save_docx(doc, original_file_name, language):
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output, f"{language}_translated_{original_file_name}"

# Function to Extract Text from Excel
def extract_text_from_xlsx(wb):
    sheet_texts = {}
    full_texts = {}

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        texts = {}
        full_text = []
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    path = f"row_{cell.row},col_{cell.column}"
                    texts[path] = [cell.value]
                    full_text.append(cell.value)
        sheet_texts[sheet_name] = texts
        full_texts[sheet_name] = "\n\n".join(full_text)

    return sheet_texts, full_texts

# Function to Apply Translated Text Back to Excel
def apply_translated_text_to_xlsx(wb, translated_dict):
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        if sheet_name in translated_dict:
            translated_sheet_dict = translated_dict[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    path = f"row_{cell.row},col_{cell.column}"
                    if path in translated_sheet_dict:
                        cell.value = translated_sheet_dict[path][0]

# Function to Extract and Translate Text from Excel
def process_xlsx(uploaded_file, target_language, progress_bar):
    file_stream = BytesIO(uploaded_file.getvalue())
    wb = load_workbook(file_stream, data_only=True)
    sheet_texts, _ = extract_text_from_xlsx(wb)
    total_sheets = len(wb.sheetnames)
    sheet_counter = 0

    for sheet_name, texts in sheet_texts.items():
        if texts:
            translated_sheet_dict = translate_text(texts, target_language)
            if translated_sheet_dict:
                apply_translated_text_to_xlsx(wb, {sheet_name: translated_sheet_dict})
        sheet_counter += 1
        # Update progress bar based on the number of processed sheets
        progress_bar.progress(sheet_counter / total_sheets, text=f"Processing sheet {sheet_counter}/{total_sheets}")

    return wb

# Function to save Excel with translations
def save_xlsx(wb, original_file_name, language):
    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output, f"{language}_translated_{original_file_name}"

# Main function
def main():
    st.set_page_config(page_title="Document Translator", page_icon=":memo:", layout='wide', initial_sidebar_state='collapsed')

    col1, col2, col3 = st.columns([1,2,1])
    with col2:
        logo_path = "bavista_logo.png" 
        st.image(logo_path, use_container_width=True) 

    st.title("Agent Philip - Document Translator")
    uploaded_file = st.file_uploader("**1. Upload your Document file**", type=["pptx", "docx", "xlsx"])

    language_option = st.radio("**2. Choose an option for language selection:**", ('Select from list', 'Enter custom language'))
    if language_option == 'Select from list':
        language = st.selectbox("Select Language you want to translate to", ['Japanese', 'Vietnamese', 'English', 'Mandarin', 'Hindi', 'Arabic', 'Spanish'])
    else:
        language = st.text_input("Enter the language you want to translate to")

    if uploaded_file and language:
        file_type = uploaded_file.name.split('.')[-1]
        if st.button(f"Translate to **{language}**", use_container_width=True, type="primary"):
            with st.spinner("🙇🏻‍♀️ Working on this task, please give it a moment..."):
                if file_type == 'pptx':
                    ppt = Presentation(uploaded_file)
                    total_slides = len(ppt.slides)
                    progress_bar = st.progress(0, text="🤔 Analyzing your slides")
                    translated_ppt = process_pptx(uploaded_file, language, progress_bar, total_slides)
                    progress_bar.progress(100, "All done ✅")
                    st.balloons()
                    translated_ppt_bytes, new_file_name = save_pptx(translated_ppt, uploaded_file.name, language)
                    st.download_button(label="💾 Download Translated PowerPoint", data=translated_ppt_bytes, file_name=new_file_name, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
                elif file_type == 'docx':
                    doc = Document(uploaded_file)
                    progress_bar = st.progress(50, text="🤔 Analyzing your documents")
                    translated_doc = process_docx(uploaded_file, language)
                    progress_bar.progress(100, "All done ✅")
                    st.balloons()
                    translated_doc_bytes, new_file_name = save_docx(translated_doc, uploaded_file.name, language)
                    st.download_button(label="💾 Download Translated Document", data=translated_doc_bytes, file_name=new_file_name, mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
                elif file_type == 'xlsx':
                    progress_bar = st.progress(0, text="🤔 Analyzing your sheets")
                    translated_wb = process_xlsx(uploaded_file, language, progress_bar)
                    progress_bar.progress(100, "All done ✅")  # Ensure the progress bar reaches 100% when done
                    st.balloons()
                    translated_xlsx_bytes, new_file_name = save_xlsx(translated_wb, uploaded_file.name, language)
                    st.download_button(label="💾 Download Translated Excel", data=translated_xlsx_bytes, file_name=new_file_name, mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True)

if __name__ == "__main__":
    main()
