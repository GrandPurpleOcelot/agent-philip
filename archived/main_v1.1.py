import streamlit as st
from pptx import Presentation
from io import BytesIO
import json
from openai import AzureOpenAI

# Set up your OpenAI API key
api_key = st.secrets["AZURE_OPENAI_API_KEY"]
azure_endpoint = st.secrets["AZURE_OPENAI_ENDPOINT"]

client = AzureOpenAI(
    api_key=api_key,  
    api_version="2024-02-01",
    azure_endpoint=azure_endpoint
)

def extract_all_text_from_ppt(file):
    ppt = Presentation(file)
    texts = {}
    full_text = []

    def extract_text_from_shapes(shapes, slide_index, path, texts, full_text):
        for index, shape in enumerate(shapes):
            current_path = f"{path},{index}"
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_text = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    shape_text.append(paragraph_text)
                    full_text.append(paragraph_text)
                if shape_text:
                    texts[current_path] = shape_text
            elif hasattr(shape, "shapes"):  # This is a group shape.
                extract_text_from_shapes(shape.shapes, slide_index, current_path, texts, full_text)

    for i, slide in enumerate(ppt.slides):
        extract_text_from_shapes(slide.shapes, i, f"{i}", texts, full_text)

    full_context = "\n\n".join(full_text)
    return texts, full_context

def extract_text_from_slide(slide):
    texts = {}
    full_text = []

    def extract_text_from_shapes(shapes, slide_index, path, texts, full_text):
        for index, shape in enumerate(shapes):
            current_path = f"{path},{index}"
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_text = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    shape_text.append(paragraph_text)
                    full_text.append(paragraph_text)
                if shape_text:
                    texts[current_path] = shape_text
            elif hasattr(shape, "shapes"):  # This is a group shape.
                extract_text_from_shapes(shape.shapes, slide_index, current_path, texts, full_text)

    extract_text_from_shapes(slide.shapes, slide.slide_id, f"{slide.slide_id}", texts, full_text)
    full_context = "\n\n".join(full_text)
    return texts, full_context

def apply_translated_text_to_slide(slide, translated_dict):
    def set_text_to_shape(shape, translated_paragraphs):
        if hasattr(shape, "text_frame"):
            for paragraph_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                if paragraph_idx < len(translated_paragraphs):
                    if paragraph.runs:
                        # Store the formatting of the first run
                        original_run = paragraph.runs[0]
                        font_size = original_run.font.size
                        font_color = original_run.font.color
                        bold = original_run.font.bold
                        italic = original_run.font.italic
                        underline = original_run.font.underline

                        # Clear existing text
                        paragraph.clear()

                        # Add new run with the translated text
                        new_run = paragraph.add_run()
                        new_run.text = translated_paragraphs[paragraph_idx]

                        # Apply the original formatting to the new run
                        if font_size:
                            new_run.font.size = font_size
                        if font_color and hasattr(font_color, 'rgb'):
                            new_run.font.color.rgb = font_color.rgb
                        new_run.font.bold = bold
                        new_run.font.italic = italic
                        new_run.font.underline = underline

    def apply_to_shapes(shapes, path):
        for index, shape in enumerate(shapes):
            current_path = f"{path},{index}"
            if current_path in translated_dict:  # Corrected from 'current_data' to 'current_path'
                set_text_to_shape(shape, translated_dict[current_path])
            elif hasattr(shape, "shapes"):
                apply_to_shapes(shape.shapes, current_path)  # Corrected recursive call

    apply_to_shapes(slide.shapes, f"{slide.slide_id}")

def translate_text(text_dict, target_language):
    # system prompt:
    prompt = f"""
        You are a professional language translator.\n
        Return a json with format similar to user's provided dictionary\n
        Translate the text to {target_language}"""
    print("------------------- System Prompt:\n", prompt)

    # user prompt:
    converted_dict = json.dumps({str(k): v for k, v in text_dict.items()})
    print("------------------- Input text dictionary:\n", converted_dict)

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            response_format={"type": "json_object"},
            messages=[{"role": "system", "content": prompt},
                      {"role": "user", "content": converted_dict}],
            temperature=0.5,
            max_tokens=4096,
        )
        print("------------------- Raw response from API:\n", response) 

        # Assuming the response format matches the input dictionary order
        translated_dict = json.loads(response.choices[0].message.content)
        return translated_dict
    except json.JSONDecodeError as e:
        print("Failed to decode JSON from response:", response)
        st.error(f"JSON decoding error: {str(e)}")
        return text_dict
    except Exception as e:
        st.error(f"Error in translation: {str(e)}")
        return text_dict

def process_pptx(file, target_language, progress_bar, total_slides):
    ppt = Presentation(file)
    for i, slide in enumerate(ppt.slides):
        slide_text_dict, _ = extract_text_from_slide(slide)
        translated_slide_dict = translate_text(slide_text_dict, target_language)
        apply_translated_text_to_slide(slide, translated_slide_dict)
        # Update progress bar based on the slide index, value between 0.0 and 1.0
        current_progress = (i + 1) / total_slides
        progress_bar.progress(current_progress, text= f"Processing slide {i}/{total_slides}")
    return ppt

def save_pptx(ppt, original_file_name, language):
    output = BytesIO()
    ppt.save(output)
    output.seek(0)
    return output, f"{language}_translated_{original_file_name}"

def main():
    # Streamlit interface
    st.set_page_config(page_title="Powerpoint Translator", page_icon=":memo:", layout='wide', initial_sidebar_state='collapsed')

    # Using columns to center the logo
    col1, col2, col3 = st.columns([1,2,1])  # Adjust the ratio as needed to center the logo
    with col2:
        logo_path = "bavista_logo.png" 
        st.image(logo_path, use_column_width=True) 

    st.title("Agent Philip - PowerPoint Translator")
    uploaded_file = st.file_uploader("Upload your PowerPoint file", type=["pptx"])
    language = st.selectbox("Select Language", ['Japanese', 'Vietnamese', 'English', 'Mandarin', 'Hindi', 'Arabic'])

    if uploaded_file and language:
        if st.button(f"Translate to **{language}**", use_container_width=True):
            with st.spinner("🙇🏻‍♀️ Working on this task, please give it a moment..."):
                ppt = Presentation(uploaded_file)
                total_slides = len(ppt.slides)
                progress_bar = st.progress(0, text="🤔 Analyzing the your slides")
                translated_ppt = process_pptx(uploaded_file, language, progress_bar, total_slides)
                progress_bar.progress(100)
                st.balloons()
                translated_ppt_bytes, new_file_name = save_pptx(translated_ppt, uploaded_file.name, language)
                st.download_button(label="Download Translated PowerPoint",
                                   data=translated_ppt_bytes,
                                   file_name=new_file_name,
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                   use_container_width=True,type="primary"
                                   )

if __name__ == "__main__":
    main()